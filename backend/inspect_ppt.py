import os
from pptx import Presentation

def inspect_details(path):
    if not os.path.exists(path):
        return
    prs = Presentation(path)
    slide = prs.slides[0]
    
    print(f"Total Shapes on Slide 1: {len(slide.shapes)}")
    for idx, shape in enumerate(slide.shapes):
        print(f"\nShape {idx+1}: {shape.name} (Type: {shape.shape_type})")
        print(f"  Position: Left={shape.left.inches:.2f}, Top={shape.top.inches:.2f}, Width={shape.width.inches:.2f}, Height={shape.height.inches:.2f}")
        
        # Check fill color safely
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                print(f"  Fill Type: {shape.fill.type}")
                # Try getting solid color
                if shape.fill.type == 1: # Solid
                    print(f"  Fill Color RGB: {shape.fill.fore_color.rgb}")
                elif shape.fill.type == 5: # Background
                    print(f"  Fill Type: Background")
        except Exception as e:
            print(f"  Fill check failed: {e}")
            
        # Check line color safely
        try:
            if hasattr(shape, 'line') and shape.line.color is not None:
                print(f"  Line Color: {shape.line.color.type}")
        except Exception as e:
            pass
        
        # Check text
        if shape.has_text_frame:
            print(f"  Text Content: {repr(shape.text_frame.text[:120])}")
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                if p.text.strip():
                    print(f"    Paragraph {p_idx+1}: {repr(p.text[:100])}")
                    if p.runs:
                        run = p.runs[0]
                        font = run.font
                        print(f"      Font: Name={font.name}, Size={font.size.pt if font.size else None}, Color={font.color.rgb if font.color.type == 1 else font.color.type if font.color else None}, Bold={font.bold}")

if __name__ == "__main__":
    inspect_details("/Users/ahzarjy/word doct/MIS.pptx")
