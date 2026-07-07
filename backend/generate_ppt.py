import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def copy_template_shapes(source_slide, target_slide):
    # Copy shapes at index 0 to 7 (background bars, university name, professor name, triangles, logo)
    for i in range(8):
        shape = source_slide.shapes[i]
        el = shape.element
        new_el = copy.deepcopy(el)
        target_slide.shapes._spTree.append(new_el)

def add_content_slide(prs, source_slide, title_text):
    # Add a blank slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Copy the template theme shapes (header, footer, logo, triangles, info)
    copy_template_shapes(source_slide, slide)
    
    # Add title for the content slide in the content area
    title_box = slide.shapes.add_textbox(Inches(0.96), Inches(1.2), Inches(11.4), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Khmer OS Muol Light"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    
    return slide

def build_appended_presentation():
    input_path = "/Users/ahzarjy/word doct/MIS.pptx"
    if not os.path.exists(input_path):
        print(f"Error: file not found at {input_path}")
        return

    # Load original presentation with the user's cover slide
    prs = Presentation(input_path)
    
    # Keep only the first slide (which is the user's cover slide)
    # This prevents duplicate slides if the script is run multiple times
    while len(prs.slides) > 1:
        # Delete slide at index 1
        rId = prs.slides._sldIdLst[1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1]
    
    cover_slide = prs.slides[0]
    
    # Text fonts
    BODY_FONT = "Kh Battambang"
    TITLE_FONT = "Khmer OS Muol Light"
    
    # Color definition
    DARK_BLUE = RGBColor(30, 41, 59)
    LIGHT_BLUE = RGBColor(14, 165, 233)
    GRAY_TEXT = RGBColor(100, 116, 139)
    WHITE = RGBColor(255, 255, 255)
    LINE_COLOR = RGBColor(99, 102, 241) # Indigo lines

    # -------------------------------------------------------------
    # SLIDE 2: Overview (бһҹбҹ’бһҗбһ¶бһ”бһҸбҹ’бһҷбһҖбһҳбҹ’бһҳ бҹЈ бһҹбҹ’бһҡбһ‘бһ¶бһ”бҹӢ)
    # -------------------------------------------------------------
    slide2 = add_content_slide(prs, cover_slide, "бһҹбҹ’бһҗбһ¶бһ”бһҸбҹ’бһҷбһҖбһҳбҹ’бһҳ бҹЈ бһҹбҹ’бһҡбһ‘бһ¶бһ”бҹӢбһ“бҹғбһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’ MIS")
    
    content_box2 = slide2.shapes.add_textbox(Inches(0.96), Inches(2.1), Inches(11.4), Inches(4.3))
    tf2 = content_box2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "бһҖбһ¶бһҡбһ”бҹӮбһ„бһ…бҹӮбһҖбһҹбһҳбһ¶бһҹбһ—бһ¶бһӮбһҮбһ¶ бҹЈ бһҹбҹ’бһҡбһ‘бһ¶бһ”бҹӢбһ’бҹҶбҹ—бҹ–"
    p.font.name = BODY_FONT
    p.font.size = Pt(15)
    p.font.bold = True
    p.space_after = Pt(10)
    
    p1 = tf2.add_paragraph()
    p1.text = "вҖў бҹЎ. Database (бһҹбҹ’бһҡбһ‘бһ¶бһ”бҹӢбһ•бҹ’бһ‘бһ»бһҖбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷ - Data Layer): бһҮбһ¶бһӮбҹ’бһҡбһ№бҹҮбһҳбһјбһӣбһҠбҹ’бһӢбһ¶бһ“бһҠбҹӮбһӣбһ•бҹ’бһ‘бһ»бһҖбһ–бҹҗбһҸбҹҢбһҳбһ¶бһ“бһ“бһ·бһ„бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһўбһ¶бһҮбһёбһңбһҖбһҳбҹ’бһҳбһ‘бһ¶бҹҶбһ„бһўбһҹбҹӢбҹ” (бһҹбҹҶбһҒбһ¶бһ“бҹӢбһ”бҹҶбһ•бһ»бһҸ)"
    p1.font.name = BODY_FONT
    p1.font.size = Pt(13)
    p1.space_after = Pt(8)
    
    p2 = tf2.add_paragraph()
    p2.text = "вҖў бҹў. Backend (бһҹбҹ’бһҡбһ‘бһ¶бһ”бҹӢбһҠбҹҶбһҺбһҫбһҡбһҖбһ¶бһҡ Logic - Application Layer): бһҠбҹҶбһҺбһҫбһҡбһҖбһ¶бһҡбһҖбҹ’бһ”бһҪбһ“бһҒбһҪбһҡбһҖбҹ’бһ”бһ¶бһӣ бһ“бһ·бһ„ API бһҹбһҳбҹ’бһҡбһ”бһҹбһҳбҹ’бһҡбһҪбһӣбһ‘бҹҶбһ“бһ¶бһҖбҹӢбһ‘бҹҶбһ“бһ„бһҡбһңбһ¶бһ„ Database бһ“бһ·бһ„ Frontendбҹ”"
    p2.font.name = BODY_FONT
    p2.font.size = Pt(13)
    p2.space_after = Pt(8)
    
    p3 = tf2.add_paragraph()
    p3.text = "вҖў бҹЈ. Frontend (бһҹбҹ’бһҡбһ‘бһ¶бһ”бҹӢбһ”бһ„бҹ’бһ бһ¶бһүбһўбҹ’бһ“бһҖбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢ - Presentation Layer): бһ…бҹҶбһҺбһ»бһ…бһ”бҹ’бһҡбһ‘бһ¶бһҖбҹӢбһҖбҹ’бһҡбһ¶бһ бҹ’бһңбһ·бһҖ (UI/UX) бһҹбһҳбҹ’бһҡбһ¶бһ”бҹӢбһұбҹ’бһҷбһўбҹ’бһ“бһҖбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢбһ‘бһјбһ‘бҹ…бһ”бһүбҹ’бһҮбһ¶ бһ“бһ·бһ„бһҳбһҫбһӣбһғбһҫбһүбһҡбһ”бһ¶бһҷбһҖбһ¶бһҡбһҺбҹҚбҹ”"
    p3.font.name = BODY_FONT
    p3.font.size = Pt(13)
    p3.space_after = Pt(18)
    
    p4 = tf2.add_paragraph()
    p4.text = "бһ…бҹҶбһҺбһ¶бҹҶбҹ– Database бһӮбһәбһҮбһ¶бһҹбһҳбһ¶бһҹбһ—бһ¶бһӮбһҹбҹҶбһҒбһ¶бһ“бҹӢбһҠбҹҶбһ”бһјбһ„бһӮбҹҒбһ”бһ„бҹ’бһўбһҹбҹӢ бһ–бҹ’бһҡбҹ„бҹҮбһ”бһҫбһӮбҹ’бһҳбһ¶бһ“бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷ бһ“бҹ„бҹҮбһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’бһ‘бһ¶бҹҶбһ„бһҳбһјбһӣбһҳбһ·бһ“бһўбһ¶бһ…бһҠбҹҶбһҺбһҫбһҡбһҖбһ¶бһҡ бһ¬бһҳбһ¶бһ“бһҸбһҳбҹ’бһӣбҹғбһўбҹ’бһңбһёбһЎбһҫбһҷбҹ”"
    p4.font.name = BODY_FONT
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = LIGHT_BLUE

    # -------------------------------------------------------------
    # SLIDE 3: Priority 1 - Database
    # -------------------------------------------------------------
    slide3 = add_content_slide(prs, cover_slide, "бһ бҹҒбһҸбһ»бһўбҹ’бһңбһёбһ”бһ¶бһ“бһҮбһ¶ Database бһҹбҹҶбһҒбһ¶бһ“бҹӢбһҮбһ¶бһ„бһӮбҹҒбһ”бһ„бҹ’бһўбһҹбҹӢ?")
    
    content_box3 = slide3.shapes.add_textbox(Inches(0.96), Inches(2.1), Inches(11.4), Inches(4.3))
    tf3 = content_box3.text_frame
    tf3.word_wrap = True
    
    points = [
        ("бҹЎ. бһҖбһ¶бһҡбһҡбһҖбҹ’бһҹбһ¶бһ‘бһ»бһҖбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһ–бһ·бһҸбһ”бҹ’бһҡбһ¶бһҖбһҠ (Data Persistence)",
         "аёЈбһ¶бһӣбҹӢбһ”бҹ’бһҡбһңбһҸбҹ’бһҸбһ·бһ“бҹғбһҖбһ¶бһҡбһӣбһҖбҹӢ бһҹбҹ’бһҸбһ»бһҖбһ‘бҹҶбһ“бһ·бһү бһ“бһ·бһ„бһӮбһҺбһ“бһёбһўбҹ’бһ“бһҖбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢ бһҸбҹ’бһҡбһјбһңбһҸбҹӮбһҡбһҖбҹ’бһҹбһ¶бһ‘бһ»бһҖбһҮбһ¶бһўбһ…бһ·бһ“бҹ’бһҸбҹ’бһҡбҹғбһҷбҹҚ бһ‘бҹ„бҹҮбһ”бһёбһҮбһ¶бһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’бһ”бһ·бһ‘ бһ¬бһӮбһ¶бҹҶбһ„бһҖбһ»бҹҶбһ–бҹ’бһҷбһјбһ‘бҹҗбһҡбһҖбҹҸбһҠбҹ„бһҷбҹ”"),
        ("бҹў. бһ—бһ¶бһ–бһҸбҹ’бһҡбһ№бһҳбһҸбҹ’бһҡбһјбһң бһ“бһ·бһ„бһҹбһ»бһңбһҸбҹ’бһҗбһ·бһ—бһ¶бһ–бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷ (Data Integrity & Security)",
         "Database бһҖбһ¶бһҡбһ–бһ¶бһҡбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһҖбһ»бҹҶбһұбҹ’бһҷбһҹбҹ’бһ‘бһҪбһ“ бһ’бһ¶бһ“бһ¶бһ‘бҹҶбһ“бһ¶бһҖбҹӢбһ‘бҹҶбһ“бһ„бһҡбһңбһ¶бһ„бһҸбһ¶бһҡбһ¶бһ„бһ•бҹ’бһҹбҹҒбһ„бҹ— (Foreign Keys) бһ“бһ·бһ„бһҖбһ¶бһҡбһ–бһ¶бһҡбһҖбһ¶бһҡбһ…бһјбһӣбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢбһҠбҹ„бһҷбһӮбҹ’бһҳбһ¶бһ“бһҖбһ¶бһҡбһўбһ“бһ»бһүбҹ’бһүбһ¶бһҸбҹ”"),
        ("бҹЈ. бһҮбһ¶бһ”бҹ’бһҡбһ—бһ–бһҹбһҳбҹ’бһҡбһ¶бһ”бҹӢбһҖбһ¶бһҡбһңбһ·бһ—бһ¶бһӮбһ–бҹҗбһҸбҹҢбһҳбһ¶бһ“ (Business Intelligence)",
         "бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһ“бҹ…бһҖбҹ’бһ“бһ»бһ„ Database бһӮбһәбһҮбһ¶бһ”бҹ’бһҡбһ—бһ–бһҸбҹӮбһҳбһҪбһҷбһӮбһҸбҹӢбһҹбһҳбҹ’бһҡбһ¶бһ”бҹӢбһҷбһҖбһҳбһҖбһӮбһҺбһ“бһ¶ бһңбһ·бһ—бһ¶бһӮ бһ“бһ·бһ„бһ‘бһ¶бһүбһҷбһҖбһҮбһ¶бһҡбһ”бһ¶бһҷбһҖбһ¶бһҡбһҺбҹҚбһӣбһҖбҹӢ бһҠбһҫбһҳбҹ’бһ”бһёбһ’бҹ’бһңбһҫбһҖбһ¶бһҡбһҹбһҳбҹ’бһҡбҹҒбһ…бһ…бһ·бһҸбҹ’бһҸбһўбһ¶бһҮбһёбһңбһҖбһҳбҹ’бһҳбҹ”")
    ]
    
    for idx, (title, desc) in enumerate(points):
        p_t = tf3.add_paragraph() if idx > 0 else tf3.paragraphs[0]
        p_t.text = title
        p_t.font.name = BODY_FONT
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = DARK_BLUE
        
        p_d = tf3.add_paragraph()
        p_d.text = desc
        p_d.font.name = BODY_FONT
        p_d.font.size = Pt(12.5)
        p_d.space_after = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 4: Priority 2 - Backend
    # -------------------------------------------------------------
    slide4 = add_content_slide(prs, cover_slide, "бһҸбһҪбһ“бһ¶бһ‘бһёбһӣбҹҶбһҠбһ¶бһ”бҹӢбһ‘бһёбһ–бһёбһҡбҹ– Backend (бһҒбһҪбһҡбһҖбҹ’бһ”бһ¶бһӣбһ“бҹғбһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’)")
    
    content_box4 = slide4.shapes.add_textbox(Inches(0.96), Inches(2.1), Inches(11.4), Inches(4.3))
    tf4 = content_box4.text_frame
    tf4.word_wrap = True
    
    p_b1 = tf4.paragraphs[0]
    p_b1.text = "вҖў бһҠбҹҶбһҺбһҫбһҡбһҖбһ¶бһҡбһҖбҹ’бһ”бһҪбһ“бһ…бҹ’бһ”бһ¶бһ”бҹӢбһўбһ¶бһҮбһёбһңбһҖбһҳбҹ’бһҳ (Business Logic)"
    p_b1.font.name = BODY_FONT
    p_b1.font.size = Pt(14)
    p_b1.font.bold = True
    p_b_d1 = tf4.add_paragraph()
    p_b_d1.text = "  - бһ•бҹ’бһ‘бҹҖбһ„бһ•бҹ’бһ‘бһ¶бһҸбҹӢбһҖбһ¶бһҡбһ”бһүбҹ’бһҮбһ¶бһ‘бһ·бһү бһӮбһҺбһ“бһ¶бһҸбһҳбҹ’бһӣбҹғбһҹбһҡбһ»бһ” бһҸбҹ’бһҡбһҪбһҸбһ–бһ·бһ“бһ·бһҸбҹ’бһҷбһҹбҹ’бһҸбһ»бһҖбһ‘бҹҶбһ“бһ·бһү бһ“бһ·бһ„бһ”бһүбҹ’бһҮбһ¶бһ‘бҹ…бһҖбһҸбҹӢбһҸбҹ’бһҡбһ¶бһҖбҹ’бһ“бһ»бһ„ Databaseбҹ”"
    p_b_d1.font.name = BODY_FONT
    p_b_d1.font.size = Pt(12.5)
    p_b_d1.space_after = Pt(10)
    
    p_b2 = tf4.add_paragraph()
    p_b2.text = "вҖў бһҹбҹ’бһ–бһ¶бһ“бһ…бһҳбҹ’бһӣбһ„бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷ (API Gateway)"
    p_b2.font.name = BODY_FONT
    p_b2.font.size = Pt(14)
    p_b2.font.bold = True
    p_b_d2 = tf4.add_paragraph()
    p_b_d2.text = "  - бһ‘бһ‘бһҪбһӣбһҹбҹҶбһҺбһҫбһ–бһё Frontend (HTTP Request) бһҡбһҪбһ…бһ‘бһ¶бһүбһҷбһҖбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһ–бһё Database бһҳбһҖбһҖбҹӮбһ…бҹ’бһ“бҹғ бһ“бһ·бһ„бһҶбҹ’бһӣбһҫбһҷбһҸбһ”бһ‘бҹ…бһңбһ·бһүбһҮбһ¶бһ‘бһҳбҹ’бһҡбһ„бҹӢ JSONбҹ”"
    p_b_d2.font.name = BODY_FONT
    p_b_d2.font.size = Pt(12.5)
    p_b_d2.space_after = Pt(10)
    
    p_b3 = tf4.add_paragraph()
    p_b3.text = "вҖў бһҹбһ»бһңбһҸбҹ’бһҗбһ·бһ—бһ¶бһ–бһҖбһҳбҹ’бһҡбһ·бһҸбһҖбһҳбҹ’бһҳбһңбһ·бһ’бһё (Application Security)"
    p_b3.font.name = BODY_FONT
    p_b3.font.size = Pt(14)
    p_b3.font.bold = True
    p_b_d3 = tf4.add_paragraph()
    p_b_d3.text = "  - бһ•бҹ’бһ‘бҹҖбһ„бһ•бҹ’бһ‘бһ¶бһҸбҹӢбһўбһҸбҹ’бһҸбһҹбһүбҹ’бһүбһ¶бһҺбһўбҹ’бһ“бһҖбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢ (JWT Authentication) бһ“бһ·бһ„бһҖбҹҶбһҺбһҸбҹӢбһҹбһ·бһ‘бҹ’бһ’бһ·бһ…бһјбһӣбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢбһӣбһҫбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһ“бһёбһҳбһҪбһҷбҹ—бҹ”"
    p_b_d3.font.name = BODY_FONT
    p_b_d3.font.size = Pt(12.5)

    # -------------------------------------------------------------
    # SLIDE 5: Priority 3 - Frontend
    # -------------------------------------------------------------
    slide5 = add_content_slide(prs, cover_slide, "бһҸбһҪбһ“бһ¶бһ‘бһёбһӣбҹҶбһҠбһ¶бһ”бҹӢбһ…бһ»бһ„бһҖбҹ’бһҡбҹ„бһҷбҹ– Frontend (бһҳбһ»бһҒбһҳбһ¶бһҸбҹӢбһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’)")
    
    content_box5 = slide5.shapes.add_textbox(Inches(0.96), Inches(2.1), Inches(11.4), Inches(4.3))
    tf5 = content_box5.text_frame
    tf5.word_wrap = True
    
    p_f1 = tf5.paragraphs[0]
    p_f1.text = "вҖў бһ…бҹҶбһҺбһ»бһ…бһ”бҹ’бһҡбһ‘бһ¶бһҖбҹӢбһўбҹ’бһ“бһҖбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢ (User Interface - UI)"
    p_f1.font.name = BODY_FONT
    p_f1.font.size = Pt(14)
    p_f1.font.bold = True
    p_f_d1 = tf5.add_paragraph()
    p_f_d1.text = "  - бһ”бһ„бҹ’бһ бһ¶бһүбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһ–бһё Database бһҖбҹ’бһ“бһ»бһ„бһ‘бһҳбҹ’бһҡбһ„бҹӢбһҠбҹӮбһӣбһ„бһ¶бһҷбһҹбҹ’бһҡбһҪбһӣбһҳбһҫбһӣ бһҠбһјбһ…бһҮбһ¶бһҸбһ¶бһҡбһ¶бһ„ бһ–бһҺбҹҢ бһҖбҹ’бһҡбһ¶бһ бҹ’бһңбһ·бһҖ бһ“бһ·бһ„бһ”бҹҠбһјбһҸбһ»бһ„бһ”бһүбҹ’бһҮбһ¶бһ•бҹ’бһҹбҹҒбһ„бҹ—бҹ”"
    p_f_d1.font.name = BODY_FONT
    p_f_d1.font.size = Pt(12.5)
    p_f_d1.space_after = Pt(10)

    p_f2 = tf5.add_paragraph()
    p_f2.text = "вҖў бһ”бһ‘бһ–бһ·бһҹбҹ„бһ’бһ“бҹҚбһўбҹ’бһ“бһҖбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢ (User Experience - UX)"
    p_f2.font.name = BODY_FONT
    p_f2.font.size = Pt(14)
    p_f2.font.bold = True
    p_f_d2 = tf5.add_paragraph()
    p_f_d2.text = "  - бһ’бһ¶бһ“бһ¶бһҖбһ¶бһҡбһҡбҹҖбһ”бһ…бҹҶбһ‘бһҳбҹ’бһҡбһ„бҹӢбһҖбһ¶бһҡбһ„бһ¶бһҡбһұбҹ’бһҷбһҳбһ¶бһ“бһ—бһ¶бһ–бһ„бһ¶бһҷбһҹбҹ’бһҡбһҪбһӣ бһҡбһ бҹҗбһҹ бһ“бһ·бһ„бһҳбһ¶бһ“бһҡбһ”бҹҖбһ”бһҡбҹҖбһ”бһҡбһҷбһҹбһҳбҹ’бһҡбһ¶бһ”бҹӢбһўбҹ’бһ“бһҖбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢбһ‘бһјбһ‘бҹ…бҹ”"
    p_f_d2.font.name = BODY_FONT
    p_f_d2.font.size = Pt(12.5)
    p_f_d2.space_after = Pt(10)

    p_f3 = tf5.add_paragraph()
    p_f3.text = "вҖў бһҖбһ¶бһҡбһ”бҹ’бһҡбһҳбһјбһӣбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһ”бһүбҹ’бһ…бһјбһӣбһҠбҹҶбһ”бһјбһ„ (Input Collection)"
    p_f3.font.name = BODY_FONT
    p_f3.font.size = Pt(14)
    p_f3.font.bold = True
    p_f_d3 = tf5.add_paragraph()
    p_f_d3.text = "  - бһ‘бһ‘бһҪбһӣбһҖбһ¶бһҡбһ”бһүбҹ’бһ…бһјбһӣбһ–бҹҗбһҸбҹҢбһҳбһ¶бһ“бһ–бһёбһўбҹ’бһ“бһҖбһ”бҹ’бһҡбһҫбһ”бҹ’бһҡбһ¶бһҹбҹӢ (бһҠбһјбһ…бһҮбһ¶бһҖбһ¶бһҡбһҖбһҸбҹӢбһҸбҹ’бһҡбһ¶бһҖбһ¶бһҡбһӣбһҖбҹӢбһҗбҹ’бһҳбһё бһ¬бһ”бһ“бҹ’бһҗбҹӮбһҳбһ‘бҹҶбһ“бһ·бһү) бһҡбһҪбһ…бһ•бҹ’бһүбһҫбһңбһ¶бһ‘бҹ…бһҖбһ¶бһ“бҹӢ Backendбҹ”"
    p_f_d3.font.name = BODY_FONT
    p_f_d3.font.size = Pt(12.5)

    # -------------------------------------------------------------
    # SLIDE 6: Database Schema / Tables
    # -------------------------------------------------------------
    slide6 = add_content_slide(prs, cover_slide, "бһҹбҹ’бһҗбһ¶бһ”бһҸбҹ’бһҷбһҖбһҳбҹ’бһҳбһҸбһ¶бһҡбһ¶бһ„бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷ (Database Tables)")
    
    content_box6 = slide6.shapes.add_textbox(Inches(0.96), Inches(2.1), Inches(11.4), Inches(4.3))
    tf6 = content_box6.text_frame
    tf6.word_wrap = True
    
    p_sch = tf6.paragraphs[0]
    p_sch.text = "бһҸбһ¶бһҡбһ¶бһ„бһҹбҹҶбһҒбһ¶бһ“бҹӢбҹ—бһ‘бһ¶бҹҶбһ„ бҹҰ бһ“бҹ…бһҖбҹ’бһ“бһ»бһ„бһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’ Sales MIS Databaseбҹ–"
    p_sch.font.name = BODY_FONT
    p_sch.font.size = Pt(14)
    p_sch.font.bold = True
    p_sch.space_after = Pt(10)
    
    sch_items = [
        ("вҖў users", "бһҡбһҖбҹ’бһҹбһ¶бһ‘бһ»бһҖбһӮбһҺбһ“бһёбһ”бһ»бһӮбҹ’бһӮбһӣбһ·бһҖ бһўбҹҠбһёбһҳбҹӮбһӣ бһӣбҹҒбһҒбһҹбһҳбҹ’бһ„бһ¶бһҸбҹӢ бһ“бһ·бһ„бһҹбһ·бһ‘бҹ’бһ’бһ·бһӮбҹ’бһҡбһ”бҹӢбһӮбҹ’бһҡбһ„ (Admin / User)бҹ”"),
        ("вҖў customers", "бһ”бҹ’бһҡбһҳбһјбһӣбһ–бҹҗбһҸбҹҢбһҳбһ¶бһ“бһ‘бҹҶбһ“бһ¶бһҖбҹӢбһ‘бҹҶбһ“бһ„бһҡбһ”бһҹбҹӢбһҖбҹ’бһҡбһ»бһҳбһ бҹҠбһ»бһ“ бһ¬бһўбһҸбһ·бһҗбһ·бһҮбһ“ бһ“бһ·бһ„бһҖбҹ’бһҡбһ»бһҳбһ…бҹҶбһҺбһ¶бһҸбҹӢбһҗбҹ’бһ“бһ¶бһҖбҹӢ (VIP, Active, Wholesale)бҹ”"),
        ("вҖў products", "бһ”бһүбҹ’бһҮбһёбһ‘бҹҶбһ“бһ·бһү бһҖбһјбһҠ (SKU) бһҸбһҳбҹ’бһӣбҹғбһӣбһҖбҹӢ бһ…бҹҶбһ“бһҪбһ“бһҹбҹ’бһҸбһ»бһҖбһҹбһҡбһ»бһ” бһ“бһ·бһ„бһҖбһҳбҹ’бһҡбһ·бһҸбһҹбҹ’бһҸбһ»бһҖбһҖбҹҶбһҺбһҸбҹӢбһҸбҹ’бһҡбһјбһңбһ”бһүбҹ’бһҮбһ¶бһ‘бһ·бһүбһ”бһ“бҹ’бһҗбҹӮбһҳ (Reorder Level)бҹ”"),
        ("вҖў orders", "бһҖбһҸбҹӢбһҸбҹ’бһҡбһ¶бһӣбҹҒбһҒбһңбһ·бһҖбҹ’бһҖбһҷбһ”бһҸбҹ’бһҡ бһҖбһ¶бһӣбһ”бһҡбһ·бһ…бҹ’бһҶбҹҒбһ‘бһӣбһҖбҹӢ бһҹбҹ’бһҗбһ¶бһ“бһ—бһ¶бһ–бһ”бһүбҹ’бһҮбһ¶бһ‘бһ·бһү (Pending, Completed, Delayed) бһ“бһ·бһ„бһҸбһҳбҹ’бһӣбҹғбһҹбһҡбһ»бһ”бҹ”"),
        ("вҖў order_items", "бһ•бҹ’бһ‘бһ»бһҖбһ–бҹҗбһҸбҹҢбһҳбһ¶бһ“бһӣбһҳбҹ’бһўбһ·бһҸбһ“бҹғбһ‘бҹҶбһ“бһ·бһүбһӣбһҖбҹӢбһ…бҹҒбһүбһҖбҹ’бһ“бһ»бһ„бһңбһ·бһҖбҹ’бһҖбһҷбһ”бһҸбҹ’бһҡбһ“бһёбһҳбһҪбһҷбҹ— (бһ…бҹҶбһ“бһҪбһ“бһӣбһҖбҹӢ бһҸбһҳбҹ’бһӣбҹғбһҜбһҖбһҸбһ¶ бһ“бһ·бһ„бһҸбһҳбҹ’бһӣбҹғбһҹбһҡбһ»бһ”бһҮбһҪбһҡ)бҹ”"),
        ("вҖў inventory_movements", "бһ”бҹ’бһҡбһңбһҸбҹ’бһҸбһ·бһ“бҹғбһҖбһ¶бһҡбһ”бһҳбҹ’бһҡбҹӮбһ”бһҳбҹ’бһҡбһҪбһӣбһҹбҹ’бһҸбһ»бһҖбһ‘бһ¶бҹҶбһ„бһўбһҹбҹӢбһҹбһҳбҹ’бһҡбһ¶бһ”бҹӢбһ’бҹ’бһңбһҫбһҹбһңбһ“бһҖбһҳбҹ’бһҳ (Audit Log) бһҠбһјбһ…бһҮбһ¶бһҖбһ¶бһҡбһӣбһҖбҹӢбһ…бҹҒбһү бһ¬бһҖбһ¶бһҡбһ”бһ“бҹ’бһҗбҹӮбһҳбһҹбҹ’бһҸбһ»бһҖбһҠбҹҶбһ”бһјбһ„бҹ”")
    ]
    
    for item_title, item_desc in sch_items:
        p_item = tf6.add_paragraph()
        p_item.text = f"{item_title}: {item_desc}"
        p_item.font.name = BODY_FONT
        p_item.font.size = Pt(12)
        p_item.space_after = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 7: Database Workflow Step-by-Step
    # -------------------------------------------------------------
    slide7 = add_content_slide(prs, cover_slide, "бһӣбҹҶбһ бһјбһҡбһҖбһ¶бһҡбһ„бһ¶бһҡбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбҹ– бһҠбҹҶбһҺбһҫбһҡбһҖбһ¶бһҡбһӣбһҖбҹӢ (Sales Workflow)")
    
    content_box7 = slide7.shapes.add_textbox(Inches(0.96), Inches(2.1), Inches(11.4), Inches(4.3))
    tf7 = content_box7.text_frame
    tf7.word_wrap = True
    
    p_flow = tf7.paragraphs[0]
    p_flow.text = "бһҹбһ„бҹ’бһңбһ¶бһҖбҹӢбһҖбһ¶бһҡбһ„бһ¶бһҡбһҡбҹҖбһ”бһ…бҹҶбһҠбҹ„бһҷбһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’ Database бһҸбһ¶бһҳбһҮбҹҶбһ бһ¶бһ“бһ“бһёбһҳбһҪбһҷбҹ—бҹ–"
    p_flow.font.name = BODY_FONT
    p_flow.font.size = Pt(14)
    p_flow.font.bold = True
    p_flow.space_after = Pt(10)
    
    steps = [
        ("бһҮбҹҶбһ бһ¶бһ“бһ‘бһё бҹЎ: бһ”бһ„бҹ’бһҖбһҫбһҸбһ”бҹ’бһҡбһҸбһ·бһ”бһҸбҹ’бһҸбһ·бһҖбһ¶бһҡбһӣбһҖбҹӢ (Create Order)",
         "бһ”бһ„бҹ’бһҖбһҫбһҸбһҮбһҪбһҡбһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһҗбҹ’бһҳбһёбһҳбһҪбһҷбһҖбҹ’бһ“бһ»бһ„бһҸбһ¶бһҡбһ¶бһ„ `orders` бһ—бҹ’бһҮбһ¶бһ”бҹӢбһ‘бҹ… `customers.id` бһ“бһ·бһ„бһҖбҹҶбһҺбһҸбҹӢбһҹбҹ’бһҗбһ¶бһ“бһ—бһ¶бһ– 'Pending'бҹ”"),
        ("бһҮбҹҶбһ бһ¶бһ“бһ‘бһё бҹў: бһ”бһүбҹ’бһ…бһјбһӣбһ”бһүбҹ’бһҮбһёбһ‘бҹҶбһ“бһ·бһүбһӣбһҳбҹ’бһўбһ·бһҸ (Insert Line Items)",
         "бһҡбһ¶бһӣбҹӢбһ‘бҹҶбһ“бһ·бһүбһҠбҹӮбһӣбһ”бһ¶бһ“бһ‘бһ·бһү бһҸбҹ’бһҡбһјбһңбһ”бһүбҹ’бһ…бһјбһӣбһ‘бҹ…бһҖбҹ’бһ“бһ»бһ„бһҸбһ¶бһҡбһ¶бһ„ `order_items` бһ—бҹ’бһҮбһ¶бһ”бҹӢбһ‘бҹ…бһҖбһ¶бһ“бҹӢ `orders.id` бһ“бһ·бһ„ `products.id` бһ“бһёбһҳбһҪбһҷбҹ—бҹ”"),
        ("бһҮбҹҶбһ бһ¶бһ“бһ‘бһё бҹЈ: бһ’бҹ’бһңбһҫбһ”бһ…бҹ’бһ…бһ»бһ”бҹ’бһ”бһ“бҹ’бһ“бһ—бһ¶бһ–бһҖбһҳбҹ’бһҡбһ·бһҸбһҹбҹ’бһҸбһ»бһҖ (Update Product Stock)",
         "бһ…бҹҶбһ“бһҪбһ“бһҹбҹ’бһҸбһ»бһҖбһҖбҹ’бһ“бһ»бһ„бһҸбһ¶бһҡбһ¶бһ„ `products` бһҸбҹ’бһҡбһјбһңбһ”бһ¶бһ“бһҠбһҖбһ…бҹҒбһүбһ‘бҹ…бһҸбһ¶бһҳбһ…бҹҶбһ“бһҪбһ“бһ‘бҹҶбһ“бһ·бһүбһҠбҹӮбһӣбһ”бһ¶бһ“бһӣбһҖбҹӢбһ…бҹҒбһүбһ“бһёбһҳбһҪбһҷбҹ—бҹ”"),
        ("бһҮбҹҶбһ бһ¶бһ“бһ‘бһё бҹӨ: бһ”бһ„бҹ’бһҖбһҫбһҸбһҖбҹҶбһҺбһҸбҹӢбһҸбҹ’бһҡбһ¶бһҹбһңбһ“бһҖбһҳбҹ’бһҳбһҹбҹ’бһҸбһ»бһҖ (Audit Inventory)",
         "бһ”бһ„бҹ’бһҖбһҫбһҸбһҖбҹҶбһҺбһҸбҹӢбһҸбҹ’бһҡбһ¶бһҗбҹ’бһҳбһёбһҳбһҪбһҷбһҖбҹ’бһ“бһ»бһ„бһҸбһ¶бһҡбһ¶бһ„ `inventory_movements` бһ”бҹ’бһҡбһ—бҹҒбһ‘ 'SALE' бһ“бһ·бһ„бһ…бҹҶбһ“бһҪбһ“бһҠбһҖ (-) бһҠбһҫбһҳбҹ’бһ”бһёбһҮбһ¶бһ—бһҹбҹ’бһҸбһ»бһҸбһ¶бһ„бһҹбһңбһ“бһҖбһҳбҹ’бһҳ (Audit Log)бҹ”")
    ]
    
    for idx, (title, desc) in enumerate(steps):
        p_s = tf7.add_paragraph()
        p_s.text = f"{title}"
        p_s.font.name = BODY_FONT
        p_s.font.size = Pt(13)
        p_s.font.bold = True
        p_s.font.color.rgb = DARK_BLUE
        
        p_sd = tf7.add_paragraph()
        p_sd.text = f"  - {desc}"
        p_sd.font.name = BODY_FONT
        p_sd.font.size = Pt(12)
        p_sd.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 8: Entity-Relationship Diagram (ER Diagram) - NEW
    # -------------------------------------------------------------
    slide8 = add_content_slide(prs, cover_slide, "бһҢбһёбһҷбҹүбһ¶бһҖбҹ’бһҡбһ¶бһҳбһ‘бҹҶбһ“бһ¶бһҖбҹӢбһ‘бҹҶбһ“бһ„бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷ (Entity-Relationship Diagram)")
    
    # 1. users table shape
    box_users = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.2), Inches(2.2), Inches(1.5))
    box_users.fill.solid()
    box_users.fill.fore_color.rgb = WHITE
    box_users.line.color.rgb = LINE_COLOR
    tf_u = box_users.text_frame
    tf_u.word_wrap = True
    p = tf_u.paragraphs[0]
    p.text = "users"
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    for field in ["id (PK)", "email", "full_name", "is_admin"]:
        p_f = tf_u.add_paragraph()
        p_f.text = field
        p_f.font.name = BODY_FONT
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = GRAY_TEXT

    # 2. customers table shape
    box_cust = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.3), Inches(2.2), Inches(1.5))
    box_cust.fill.solid()
    box_cust.fill.fore_color.rgb = WHITE
    box_cust.line.color.rgb = LINE_COLOR
    tf_c = box_cust.text_frame
    tf_c.word_wrap = True
    p = tf_c.paragraphs[0]
    p.text = "customers"
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    for field in ["id (PK)", "customer_code", "name", "email", "segment"]:
        p_f = tf_c.add_paragraph()
        p_f.text = field
        p_f.font.name = BODY_FONT
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = GRAY_TEXT

    # 3. inventory_movements table shape (Middle Top)
    box_inv = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(2.2), Inches(2.7), Inches(1.5))
    box_inv.fill.solid()
    box_inv.fill.fore_color.rgb = WHITE
    box_inv.line.color.rgb = LINE_COLOR
    tf_i = box_inv.text_frame
    tf_i.word_wrap = True
    p = tf_i.paragraphs[0]
    p.text = "inventory_movements"
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    for field in ["id (PK)", "product_id (FK)", "movement_type", "quantity", "created_by (FK)"]:
        p_f = tf_i.add_paragraph()
        p_f.text = field
        p_f.font.name = BODY_FONT
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = GRAY_TEXT

    # 4. orders table shape (Middle Bottom)
    box_ord = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(4.3), Inches(2.7), Inches(1.5))
    box_ord.fill.solid()
    box_ord.fill.fore_color.rgb = WHITE
    box_ord.line.color.rgb = LINE_COLOR
    tf_o = box_ord.text_frame
    tf_o.word_wrap = True
    p = tf_o.paragraphs[0]
    p.text = "orders"
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    for field in ["id (PK)", "order_number", "customer_id (FK)", "total", "ordered_at"]:
        p_f = tf_o.add_paragraph()
        p_f.text = field
        p_f.font.name = BODY_FONT
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = GRAY_TEXT

    # 5. products table shape (Right Top)
    box_prod = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(2.2), Inches(2.2), Inches(1.5))
    box_prod.fill.solid()
    box_prod.fill.fore_color.rgb = WHITE
    box_prod.line.color.rgb = LINE_COLOR
    tf_p = box_prod.text_frame
    tf_p.word_wrap = True
    p = tf_p.paragraphs[0]
    p.text = "products"
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    for field in ["id (PK)", "product_code", "sku", "name", "price", "stock"]:
        p_f = tf_p.add_paragraph()
        p_f.text = field
        p_f.font.name = BODY_FONT
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = GRAY_TEXT

    # 6. order_items table shape (Right Bottom)
    box_items = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(4.3), Inches(2.2), Inches(1.5))
    box_items.fill.solid()
    box_items.fill.fore_color.rgb = WHITE
    box_items.line.color.rgb = LINE_COLOR
    tf_it = box_items.text_frame
    tf_it.word_wrap = True
    p = tf_it.paragraphs[0]
    p.text = "order_items"
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    for field in ["id (PK)", "order_id (FK)", "product_id (FK)", "quantity", "unit_price"]:
        p_f = tf_it.add_paragraph()
        p_f.text = field
        p_f.font.name = BODY_FONT
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = GRAY_TEXT

    # Drawing connector lines between tables
    # 1. users -> inventory_movements (Horiz)
    line1 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.4), Inches(2.95), Inches(1.4), Inches(0.02))
    line1.fill.solid()
    line1.fill.fore_color.rgb = LINE_COLOR
    line1.line.fill.background()

    # 2. customers -> orders (Horiz)
    line2 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.4), Inches(5.05), Inches(1.4), Inches(0.02))
    line2.fill.solid()
    line2.fill.fore_color.rgb = LINE_COLOR
    line2.line.fill.background()

    # 3. products -> inventory_movements (Horiz)
    line3 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(2.95), Inches(1.4), Inches(0.02))
    line3.fill.solid()
    line3.fill.fore_color.rgb = LINE_COLOR
    line3.line.fill.background()

    # 4. products -> order_items (Vert)
    line4 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.0), Inches(3.7), Inches(0.02), Inches(0.6))
    line4.fill.solid()
    line4.fill.fore_color.rgb = LINE_COLOR
    line4.line.fill.background()

    # 5. orders -> order_items (Horiz)
    line5 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(5.05), Inches(1.4), Inches(0.02))
    line5.fill.solid()
    line5.fill.fore_color.rgb = LINE_COLOR
    line5.line.fill.background()

    # Note text box for explanations
    note_box = slide8.shapes.add_textbox(Inches(1.2), Inches(6.0), Inches(9.9), Inches(0.5))
    note_tf = note_box.text_frame
    note_tf.word_wrap = True
    note_p = note_tf.paragraphs[0]
    note_p.text = "бһ‘бҹҶбһ“бһ¶бһҖбҹӢбһ‘бҹҶбһ“бһ„бһҡбһңбһ¶бһ„бһҸбһ¶бһҡбһ¶бһ„ (Relationships)бҹ–\nвҖў PK: Primary Key (бһҹбҹ„бһ…бһҳбҹ’бһ”бһ„) | FK: Foreign Key (бһҹбҹ„бһҖбҹ’бһҡбҹ…бһҠбҹӮбһӣбһ—бҹ’бһҮбһ¶бһ”бҹӢбһ‘бҹҶбһ“бһ¶бһҖбҹӢбһ‘бҹҶбһ“бһ„бһҡбһңбһ¶бһ„бһҸбһ¶бһҡбһ¶бһ„бһҳбһҪбһҷбһ‘бҹ…бһҸбһ¶бһҡбһ¶бһ„бһҳбһҪбһҷбһ‘бҹҖбһҸ)"
    note_p.font.name = BODY_FONT
    note_p.font.size = Pt(11)
    note_p.font.color.rgb = GRAY_TEXT

    # -------------------------------------------------------------
    # SLIDE 9: Conclusion
    # -------------------------------------------------------------
    slide9 = add_content_slide(prs, cover_slide, "бһҹбһ“бҹ’бһ“бһ·бһҠбҹ’бһӢбһ¶бһ“бһҡбһҪбһҳ (Conclusion)")
    
    content_box9 = slide9.shapes.add_textbox(Inches(0.96), Inches(2.1), Inches(11.4), Inches(4.3))
    tf9 = content_box9.text_frame
    tf9.word_wrap = True
    
    p_c1 = tf9.paragraphs[0]
    p_c1.text = "вҖў Database (бһҳбһјбһӣбһҠбҹ’бһӢбһ¶бһ“бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷ) бһӮбһәбһҮбһ¶ бһӮбҹ’бһҡбһ№бҹҮбһҹбҹ’бһҗбһ¶бһ“ (Foundation) бһҠбҹҸбһҡбһ№бһ„бһҳбһ¶бҹҶ бһҠбҹӮбһӣбһҡбһҖбҹ’бһҹбһ¶бһ‘бһ»бһҖбһ–бҹҗбһҸбҹҢбһҳбһ¶бһ“бһ“бһ·бһ„бһ‘бҹ’бһҡбһ–бҹ’бһҷбһҹбһҳбҹ’бһ”бһҸбҹ’бһҸбһ·бһҹбҹ’бһ“бһјбһӣбҹ”"
    p_c1.font.name = BODY_FONT
    p_c1.font.size = Pt(14)
    p_c1.space_after = Pt(12)

    p_c2 = tf9.add_paragraph()
    p_c2.text = "вҖў Backend (бһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’бһҠбҹҶбһҺбһҫбһҡбһҖбһ¶бһҡбһҒбһ¶бһ„бһҖбҹ’бһҡбҹ„бһҷ) бһӮбһәбһҮбһ¶ бһҳбҹүбһ¶бһҹбҹҠбһёбһ“ (Engine) бһҠбҹҸбһҳбһ¶бһ“бһҗбһ¶бһҳбһ–бһӣ бһҠбҹӮбһӣбһӮбһҺбһ“бһ¶ бһ“бһ·бһ„бһ…бһҳбҹ’бһӣбһ„бһ‘бһ·бһ“бҹ’бһ“бһ“бҹҗбһҷбһ‘бҹ…бһҳбһҖбҹ”"
    p_c2.font.name = BODY_FONT
    p_c2.font.size = Pt(14)
    p_c2.space_after = Pt(12)

    p_c3 = tf9.add_paragraph()
    p_c3.text = "вҖў Frontend (бһ…бҹҶбһҺбһ»бһ…бһ”бҹ’бһҡбһ‘бһ¶бһҖбҹӢбһҒбһ¶бһ„бһҳбһ»бһҒ) бһӮбһәбһҮбһ¶ бһҳбһ»бһҒбһҳбһ¶бһҸбҹӢ (Interface) бһҠбҹҸбһҹбҹ’бһҡбһҹбҹӢбһҹбҹ’бһўбһ¶бһҸ бһҠбҹӮбһӣбһҮбһҪбһҷбһұбҹ’бһҷбһҳбһ“бһ»бһҹбҹ’бһҹбһ„бһ¶бһҷбһҹбҹ’бһҡбһҪбһӣбһ’бҹ’бһңбһҫбһҖбһ¶бһҡбһҮбһ¶бһҳбһҪбһҷбһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’бҹ”"
    p_c3.font.name = BODY_FONT
    p_c3.font.size = Pt(14)
    p_c3.space_after = Pt(20)

    p_c4 = tf9.add_paragraph()
    p_c4.text = "бһҖбһ¶бһҡбһҡбһҪбһҳбһ”бһүбҹ’бһ…бһјбһӣбһӮбҹ’бһ“бһ¶бһҷбҹүбһ¶бһ„бһӣбҹ’бһўбһҘбһҸбһҒбҹ’бһ…бҹ„бҹҮбһҡбһңбһ¶бһ„бһҹбһҳбһ¶бһҹбһ—бһ¶бһӮбһ‘бһ¶бҹҶбһ„бһ”бһёбһ“бҹҒбҹҮ бһ“бһ№бһ„бһ”бһ„бҹ’бһҖбһҫбһҸбһ”бһ¶бһ“бһҮбһ¶бһ”бҹ’бһҡбһ–бҹҗбһ“бҹ’бһ’ MIS бһҳбһҪбһҷбһҠбҹҸбһҮбҹ„бһӮбһҮбҹҗбһҷ бһ“бһ·бһ„бһҳбһ¶бһ“бһ”бҹ’бһҡбһҹбһ·бһ‘бҹ’бһ’бһ—бһ¶бһ–бһҒбҹ’бһ–бһҹбҹӢбһҹбһҳбҹ’бһҡбһ¶бһ”бҹӢбһўбһ¶бһҮбһёбһңбһҖбһҳбҹ’бһҳбҹ”"
    p_c4.font.name = BODY_FONT
    p_c4.font.size = Pt(14.5)
    p_c4.font.bold = True
    p_c4.font.color.rgb = LIGHT_BLUE
    p_c4.alignment = PP_ALIGN.CENTER

    # Save back to the user's original presentation location
    prs.save(input_path)
    print(f"Successfully appended slides to original presentation at: {input_path}")

if __name__ == "__main__":
    build_appended_presentation()
